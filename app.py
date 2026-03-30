import os
import asyncio
import json
import re
import uuid
import time
import pdfplumber
from pathlib import Path
from typing import Optional
from datetime import datetime, timedelta

from fastapi import FastAPI, UploadFile, File, Form, Request, HTTPException, WebSocket, WebSocketDisconnect
from fastapi.responses import HTMLResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates

try:
    import anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False

from fastapi.middleware.cors import CORSMiddleware

app = FastAPI(title="RFP AI Analyzer", version="2.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])
app.mount("/static", StaticFiles(directory="static"), name="static")
templates = Jinja2Templates(directory="templates")

import db

UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

DATA_DIR = Path("data")
DATA_DIR.mkdir(exist_ok=True)


class ConnectionManager:
    def __init__(self):
        self.active: list[dict] = []  # [{ws, username}]

    async def connect(self, ws: WebSocket, username: str):
        await ws.accept()
        self.active.append({"ws": ws, "username": username})
        await self.broadcast_users()

    def disconnect(self, ws: WebSocket):
        self.active = [c for c in self.active if c["ws"] != ws]

    async def broadcast_users(self):
        users = list({c["username"] for c in self.active})
        await self.broadcast({"type": "users", "users": users, "count": len(users)})

    async def broadcast(self, data: dict):
        dead = []
        for c in self.active:
            try:
                await c["ws"].send_json(data)
            except Exception:
                dead.append(c)
        for c in dead:
            self.active.remove(c)

    async def notify(self, action: str, detail: str, username: str = ""):
        await self.broadcast({
            "type": "activity",
            "action": action,
            "detail": detail,
            "username": username,
            "time": datetime.now().strftime("%H:%M:%S"),
        })


ws_manager = ConnectionManager()


def log_activity(action: str, detail: str = ""):
    db.insert_activity(datetime.now().strftime("%H:%M:%S"), action, detail)


def record_history(rfp_id: str, step: str, result: str):
    if not rfp_id:
        return
    ver = db.count_history_step(rfp_id, step) + 1
    db.insert_history(rfp_id, step, ver, datetime.now().strftime("%Y-%m-%d %H:%M:%S"), result)


def get_anthropic_client():
    api_key = os.environ.get("ANTHROPIC_API_KEY", "")
    if ANTHROPIC_AVAILABLE and api_key:
        return anthropic.Anthropic(api_key=api_key)
    return None


def call_ai(system_prompt: str, user_prompt: str, mock_type: str = "") -> str:
    client = get_anthropic_client()
    if client:
        try:
            message = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=8192,
                system=system_prompt,
                messages=[{"role": "user", "content": user_prompt}],
            )
            return message.content[0].text
        except Exception as e:
            return f"[AI API 오류: {e}]\n\n" + generate_mock_response(mock_type)
    return generate_mock_response(mock_type)


def generate_mock_response(mock_type: str) -> str:
    if mock_type == "analyze":
        return json.dumps({
            "summary": "RFP 문서 자동 분석 결과",
            "requirements": [
                {"id": "REQ-001", "category": "기능", "description": "시스템 로그인 및 사용자 인증 기능", "priority": "필수", "risk": "낮음"},
                {"id": "REQ-002", "category": "기능", "description": "데이터 실시간 대시보드 구현", "priority": "필수", "risk": "중간"},
                {"id": "REQ-003", "category": "성능", "description": "동시 사용자 1,000명 이상 지원", "priority": "필수", "risk": "높음"},
                {"id": "REQ-004", "category": "보안", "description": "개인정보 암호화 저장 및 전송", "priority": "필수", "risk": "높음"},
                {"id": "REQ-005", "category": "UI/UX", "description": "반응형 웹 디자인 적용", "priority": "선택", "risk": "낮음"},
                {"id": "REQ-006", "category": "연동", "description": "기존 ERP 시스템과 API 연동", "priority": "필수", "risk": "높음"},
                {"id": "REQ-007", "category": "운영", "description": "시스템 운영 매뉴얼 및 교육 제공", "priority": "선택", "risk": "낮음"},
            ],
            "evaluation_criteria": [
                {"criteria": "기술 이해도", "weight": "30%", "description": "RFP 요구사항에 대한 기술적 이해와 해결 방안"},
                {"criteria": "수행 경험", "weight": "25%", "description": "유사 프로젝트 수행 실적 및 레퍼런스"},
                {"criteria": "가격 적정성", "weight": "20%", "description": "제안 금액의 적정성 및 비용 효율성"},
                {"criteria": "프로젝트 관리", "weight": "15%", "description": "일정 관리, 리스크 관리 방안"},
                {"criteria": "기술 지원", "weight": "10%", "description": "유지보수 및 기술 지원 계획"},
            ],
            "risks": [
                {"risk": "일정 지연 리스크", "level": "높음", "description": "ERP 연동 복잡도로 인한 일정 초과 가능성"},
                {"risk": "성능 리스크", "level": "중간", "description": "대규모 동시 접속 처리 시 병목 가능성"},
                {"risk": "보안 리스크", "level": "높음", "description": "개인정보 관련 규제 준수 필요"},
            ],
        }, ensure_ascii=False, indent=2)
    elif mock_type == "pattern":
        return json.dumps({
            "industry_analysis": "IT/디지털 전환 프로젝트",
            "winning_patterns": [
                {"pattern": "기술 역량 강조", "description": "클라우드 네이티브 아키텍처 경험과 마이크로서비스 전환 실적을 전면에 배치", "confidence": "높음"},
                {"pattern": "ROI 수치화", "description": "도입 후 예상 비용 절감률(30~40%)과 생산성 향상 지표를 구체적으로 제시", "confidence": "높음"},
                {"pattern": "리스크 선제 대응", "description": "예상 리스크와 대응 방안을 별도 섹션으로 구성하여 신뢰도 확보", "confidence": "중간"},
                {"pattern": "단계적 접근", "description": "Big Bang이 아닌 Phase별 접근으로 리스크 최소화 전략 제시", "confidence": "높음"},
            ],
            "style_recommendations": [
                "공공기관: 안정성, 준법성, 보안 인증 강조",
                "대기업: ROI, 확장성, 글로벌 레퍼런스 강조",
                "중견기업: 비용 효율성, 빠른 도입, 맞춤 지원 강조",
            ],
            "differentiation_tips": [
                "경쟁사 대비 차별화 포인트: AI/자동화 역량",
                "유사 프로젝트 성공 사례를 구체적 수치와 함께 제시",
                "고객 맞춤형 PoC(Proof of Concept) 제안으로 신뢰 확보",
            ],
        }, ensure_ascii=False, indent=2)
    elif mock_type == "proposal":
        return json.dumps({
            "title": "디지털 전환 플랫폼 구축 제안서",
            "table_of_contents": [
                "1. 제안 개요",
                "  1.1 제안 배경 및 목적",
                "  1.2 프로젝트 범위",
                "2. 현황 분석",
                "  2.1 고객 환경 분석",
                "  2.2 개선 방향",
                "3. 제안 솔루션",
                "  3.1 시스템 아키텍처",
                "  3.2 핵심 기능 상세",
                "  3.3 기술 스택",
                "4. 수행 방안",
                "  4.1 프로젝트 추진 체계",
                "  4.2 일정 계획",
                "  4.3 품질 관리 방안",
                "5. 수행 실적",
                "  5.1 유사 프로젝트 레퍼런스",
                "  5.2 투입 인력 현황",
                "6. 유지보수 및 지원",
                "7. 투자 비용",
            ],
            "sections": {
                "1. 제안 개요": "본 제안서는 귀사의 디지털 전환을 위한 통합 플랫폼 구축 방안을 제시합니다. RFP에서 요구하신 핵심 요구사항을 분석하여, 안정적이고 확장 가능한 시스템 구축을 통해 업무 효율성 향상과 비용 절감을 실현하고자 합니다.",
                "3. 제안 솔루션": "마이크로서비스 아키텍처 기반의 클라우드 네이티브 플랫폼을 제안합니다. Kubernetes 기반 컨테이너 오케스트레이션으로 자동 스케일링을 지원하며, API Gateway를 통한 서비스 통합으로 기존 ERP 시스템과의 원활한 연동을 보장합니다.",
                "4. 수행 방안": "3단계 접근법(Phase 1: 기반 구축 2개월, Phase 2: 핵심 기능 개발 3개월, Phase 3: 통합 테스트 및 안정화 1개월)으로 총 6개월 내 안정적인 시스템 구축을 완료합니다.",
            },
        }, ensure_ascii=False, indent=2)
    elif mock_type == "review":
        return json.dumps({
            "overall_score": 78,
            "grade": "B+",
            "review_items": [
                {"category": "요구사항 충족", "score": 85, "status": "양호", "comment": "핵심 요구사항 7개 중 6개 충족, REQ-006 ERP 연동 상세 방안 보완 필요"},
                {"category": "논리 일관성", "score": 80, "status": "양호", "comment": "전체적으로 일관적이나, 일정 계획과 투입 인력 간 불일치 발견"},
                {"category": "차별화 요소", "score": 65, "status": "보완필요", "comment": "경쟁사 대비 뚜렷한 차별점 부족, AI/자동화 역량 강조 권장"},
                {"category": "가격 경쟁력", "score": 75, "status": "보통", "comment": "시장 평균 대비 적정 수준이나, 가치 기반 설명 보강 필요"},
                {"category": "문서 완성도", "score": 82, "status": "양호", "comment": "구성 체계적이나, 시각 자료(다이어그램) 추가 권장"},
            ],
            "missing_requirements": ["ERP 연동 상세 인터페이스 명세", "재해복구(DR) 방안", "데이터 마이그레이션 계획"],
            "logic_issues": ["3.2절 '6개월 완료' vs 4.2절 '8개월 일정표' 불일치", "투입 인력 10명 제안 vs 유사 프로젝트 15명 투입 사례 괴리"],
            "improvement_suggestions": ["ERP 연동 아키텍처 다이어그램 추가", "Phase별 마일스톤과 산출물 명확화", "경쟁사 대비 기술 우위 표 작성", "고객 인터뷰 기반 Pain Point 해결 사례 추가"],
        }, ensure_ascii=False, indent=2)
    elif mock_type == "strategy":
        return json.dumps({
            "recommendation": "GO",
            "confidence": "75%",
            "analysis": {
                "market_fit": {"score": 80, "comment": "당사 핵심 역량과 높은 부합도"},
                "competition": {"score": 65, "comment": "3~4개 경쟁사 예상, 중상위 경쟁 강도"},
                "profitability": {"score": 70, "comment": "예상 마진율 15~20%, 양호"},
                "strategic_value": {"score": 85, "comment": "해당 산업군 레퍼런스 확보 시 후속 사업 기회"},
                "resource_availability": {"score": 75, "comment": "핵심 인력 확보 가능, 일부 외부 충원 필요"},
            },
            "key_factors": [
                "유사 프로젝트 3건 수행 실적 보유",
                "해당 고객사와 기존 유지보수 관계",
                "제안 마감까지 3주, 촉박한 일정",
                "핵심 PM 다른 프로젝트 투입 중",
                "요구 기술(SAP 연동) 내부 경험 부족",
            ],
            "win_strategy": [
                "기존 고객 관계 활용하여 사전 미팅 추진",
                "SAP 연동 전문 파트너사 확보",
                "PoC 제안으로 기술 역량 증명",
                "공격적 가격 전략 (초기 수익보다 레퍼런스 확보 우선)",
            ],
        }, ensure_ascii=False, indent=2)
    elif mock_type == "knowledge":
        return json.dumps({
            "saved": True,
            "message": "지식 자산이 성공적으로 저장되었습니다.",
            "categories": ["기술 제안", "프로젝트 관리", "보안/인증"],
            "reusable_count": 5,
        }, ensure_ascii=False, indent=2)
    elif mock_type == "estimate":
        return json.dumps({
            "project_name": "클라우드 플랫폼 구축",
            "total_cost": "5억 2,000만원", "total_cost_number": 520000000,
            "duration_months": 6, "summary": "RFP 요구사항 기반 6개월 프로젝트 견적입니다.",
            "categories": [
                {"name": "인건비", "subtotal": "3억 9,600만원", "subtotal_number": 396000000, "ratio": "76%", "items": [
                    {"role": "PM", "grade": "특급", "count": 1, "months": 6, "unit_cost": "450만원", "cost": "2,700만원", "reason": "프로젝트 총괄 관리"},
                    {"role": "아키텍트", "grade": "특급", "count": 1, "months": 6, "unit_cost": "450만원", "cost": "2,700만원", "reason": "시스템 아키텍처 설계"},
                    {"role": "백엔드 개발", "grade": "고급", "count": 3, "months": 6, "unit_cost": "380만원", "cost": "6,840만원", "reason": "API 및 핵심 로직 개발"},
                    {"role": "프론트엔드 개발", "grade": "고급", "count": 2, "months": 6, "unit_cost": "380만원", "cost": "4,560만원", "reason": "UI/UX 구현"},
                    {"role": "QA", "grade": "중급", "count": 2, "months": 4, "unit_cost": "310만원", "cost": "2,480만원", "reason": "통합 테스트"},
                ]},
                {"name": "SW 라이선스", "subtotal": "5,400만원", "subtotal_number": 54000000, "ratio": "10%", "items": [
                    {"item": "클라우드 서비스 (AWS)", "cost": "3,600만원", "reason": "EKS, RDS, ElastiCache 등 운영 인프라"},
                    {"item": "모니터링 (Datadog)", "cost": "1,800만원", "reason": "APM, 로그 관리, 알림"},
                ]},
                {"name": "HW/인프라", "subtotal": "4,000만원", "subtotal_number": 40000000, "ratio": "8%", "items": [
                    {"item": "개발/스테이징 서버", "cost": "2,400만원", "reason": "개발 및 테스트 환경 구축"},
                    {"item": "네트워크/보안 장비", "cost": "1,600만원", "reason": "방화벽, VPN, SSL 인증서"},
                ]},
                {"name": "기타 경비", "subtotal": "3,000만원", "subtotal_number": 30000000, "ratio": "6%", "items": [
                    {"item": "프로젝트 관리비", "cost": "2,000만원", "reason": "회의, 출장, 문서화"},
                    {"item": "교육/인수인계", "cost": "1,000만원", "reason": "운영팀 교육 및 매뉴얼 작성"},
                ]},
            ],
            "risks": [
                {"risk": "요구사항 변경", "impact": "10~15% 추가 비용", "mitigation": "변경관리 프로세스 수립"},
                {"risk": "인력 수급 지연", "impact": "1~2개월 일정 지연", "mitigation": "핵심 인력 사전 확보"},
            ],
            "assumptions": ["SW기술자 노임단가 2024년 기준 적용", "클라우드 인프라 종량제 기준", "6개월 고정 기간 산정"],
            "notes": "실제 계약 시 요구사항 확정 후 ±10% 조정 가능",
        }, ensure_ascii=False, indent=2)
    return "분석이 완료되었습니다."


def clean_text(text: str) -> str:
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)


def extract_pdf_text(filepath: str) -> str:
    text = ""
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        text = f"[PDF 파싱 오류: {e}]"
    return clean_text(text)


# ─── Routes ───

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    return templates.TemplateResponse("index.html", {"request": request})


# ─── Dashboard API ───

@app.get("/api/dashboard")
async def dashboard():
    rfps = db.list_rfps()
    pipelines = db.list_pipelines()
    return JSONResponse({
        "rfp_count": len(rfps),
        "proposal_count": db.count_proposals(),
        "knowledge_count": db.count_knowledge(),
        "pipeline_count": db.count_pipelines(),
        "recent_rfps": [
            {"id": v["id"], "filename": v["filename"], "text_length": v["text_length"]}
            for v in rfps[:5]
        ],
        "activity_log": db.get_recent_activities(10),
        "pipeline_status": {
            rfp_id: {
                "steps": list(data.get("completed_steps", {}).keys()),
                "total": len(data.get("completed_steps", {})),
            }
            for rfp_id, data in pipelines.items()
        },
    })


# ─── Upload ───

@app.post("/api/upload-rfp")
async def upload_rfp(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "PDF 파일만 업로드 가능합니다.")

    rfp_id = str(uuid.uuid4())[:8]
    filepath = UPLOAD_DIR / f"{rfp_id}_{file.filename}"
    content = await file.read()
    with open(filepath, "wb") as f:
        f.write(content)

    text = extract_pdf_text(str(filepath))
    db.insert_rfp(rfp_id, file.filename, str(filepath), len(text), datetime.now().isoformat())
    db.upsert_pipeline(rfp_id, {}, {})
    log_activity("RFP 업로드", f"{file.filename} ({len(text):,}자)")
    return {"rfp_id": rfp_id, "filename": file.filename, "text_length": len(text), "preview": text[:500]}


# ─── 1. RFP 자동 구조화 ───

@app.post("/api/analyze-rfp")
async def analyze_rfp(rfp_id: str = Form(...)):
    rfp = db.get_rfp_meta(rfp_id)
    if not rfp:
        raise HTTPException(404, "RFP를 찾을 수 없습니다.")
    rfp["text"] = extract_pdf_text(rfp["filepath"]) if rfp.get("filepath") and Path(rfp["filepath"]).exists() else ""

    system = """당신은 RFP(제안요청서) 분석 전문가입니다.
RFP 원문에 실제로 기재된 내용만 추출하세요. 원문에 없는 내용을 임의로 만들지 마세요.

중요 규칙:
- id: RFP 원문에 요구사항 번호가 있으면 그대로 사용하고, 없으면 순번(1, 2, 3...)으로 표기
- description: 반드시 RFP 원문의 문장을 인용하거나 요약. 원문에 없는 요구사항을 절대 만들지 말 것
- evaluation_criteria: RFP에 평가 기준이 명시되어 있을 때만 추출. 없으면 빈 배열
- risks: RFP 내용에서 추론 가능한 리스크만 표기

반드시 아래 JSON 형식으로만 반환. 마크다운 코드블록 없이 순수 JSON만 출력:
{"summary":"RFP 핵심 요약","requirements":[{"id":"원문번호 또는 순번","category":"분류","description":"원문 기반 설명","priority":"필수/선택","risk":"높음/중간/낮음"}],"evaluation_criteria":[{"criteria":"기준","weight":"배점","description":"설명"}],"risks":[{"risk":"리스크명","level":"높음/중간/낮음","description":"설명"}]}"""

    result = call_ai(system, f"다음 RFP를 구조화해주세요:\n\n{rfp['text'][:8000]}", mock_type="analyze")
    db.update_pipeline_step(rfp_id, "analyze", result)
    record_history(rfp_id, "analyze", result)
    log_activity("RFP 구조화", rfp["filename"])
    await ws_manager.notify("RFP 구조화 완료", rfp["filename"])
    return JSONResponse({"rfp_id": rfp_id, "analysis": result})


# ─── 2. Winning 패턴 ───

@app.post("/api/winning-pattern")
async def winning_pattern(
    rfp_id: str = Form(None),
    industry: str = Form("IT"),
    customer_type: str = Form("대기업"),
):
    rfp_text = ""
    if rfp_id and db.rfp_exists(rfp_id):
        meta = db.get_rfp_meta(rfp_id)
        rfp_text = extract_pdf_text(meta["filepath"])[:4000] if meta and Path(meta["filepath"]).exists() else ""

    system = """당신은 제안서 수주 전략 전문가입니다.
고객 유형과 산업 분야를 고려하여 Winning Proposal 패턴을 반드시 아래 JSON 형식으로만 반환하세요. 마크다운 코드블록 없이 순수 JSON만 출력하세요:
{"industry_analysis":"산업분석","winning_patterns":[{"pattern":"패턴명","description":"설명","confidence":"높음/중간/낮음"}],"style_recommendations":["추천1"],"differentiation_tips":["전략1"]}"""

    user_msg = f"산업: {industry}\n고객 유형: {customer_type}"
    if rfp_text:
        user_msg += f"\n\nRFP 내용:\n{rfp_text}"

    result = call_ai(system, user_msg, mock_type="pattern")
    if rfp_id:
        db.update_pipeline_step(rfp_id, "pattern", result)
    record_history(rfp_id, "pattern", result)
    log_activity("패턴 분석", f"{industry} / {customer_type}")
    return JSONResponse({"analysis": result})


# ─── 3. 제안서 초안 ───

@app.post("/api/generate-proposal")
async def generate_proposal(
    rfp_id: str = Form(None),
    company_info: str = Form(""),
    references: str = Form(""),
):
    rfp_text = ""
    if rfp_id and db.rfp_exists(rfp_id):
        meta = db.get_rfp_meta(rfp_id)
        rfp_text = extract_pdf_text(meta["filepath"])[:6000] if meta and Path(meta["filepath"]).exists() else ""

    system = """당신은 20년 경력의 제안서 작성 전문 컨설턴트입니다.
RFP 요구사항, 회사 정보, 레퍼런스를 결합하여 **바로 제출 가능한 수준**의 제안서 초안을 작성하세요.

## 작성 원칙
1. 각 섹션은 최소 200자 이상, 구체적이고 설득력 있는 내용으로 작성
2. RFP의 평가 기준에 직접 대응하는 내용을 포함
3. 정량적 수치(기간, 인원, 비용절감률 등)를 적극 활용
4. 고객의 Pain Point를 정확히 짚고 해결 방안을 제시
5. 한국어 비즈니스 공식 문체 사용 (존칭, 경어체)
6. 차별화 포인트를 각 섹션에 자연스럽게 녹여 서술

## 필수 포함 섹션
- 제안 개요: 배경, 목적, 기대효과를 명확히 서술
- 현황 분석: 고객 환경 분석, AS-IS/TO-BE 비교
- 제안 솔루션: 시스템 아키텍처, 핵심 기능 상세, 기술 스택과 선정 근거
- 수행 방안: 추진 체계, WBS 기반 일정 계획, 단계별 산출물, 품질/리스크 관리
- 수행 실적: 유사 프로젝트 레퍼런스 (정량적 성과 포함)
- 투입 인력: 핵심 인력 역할과 경력 요약
- 유지보수: 하자보수, SLA, 기술지원 체계
- 투자 비용: 비용 구조 개요

반드시 아래 JSON 형식으로만 반환하세요. 마크다운 코드블록 없이 순수 JSON만 출력:
{"title":"제안서 제목","table_of_contents":["1. 제안 개요","  1.1 제안 배경 및 목적","  1.2 프로젝트 범위 및 기대효과","2. 현황 분석","  2.1 고객 환경 분석","  2.2 개선 방향 (AS-IS/TO-BE)","3. 제안 솔루션","  3.1 시스템 아키텍처","  3.2 핵심 기능 상세","  3.3 기술 스택 및 선정 근거","4. 수행 방안","  4.1 프로젝트 추진 체계","  4.2 일정 계획 (WBS)","  4.3 품질 및 리스크 관리","5. 수행 실적 및 투입 인력","  5.1 유사 프로젝트 레퍼런스","  5.2 핵심 투입 인력","6. 유지보수 및 기술 지원","7. 투자 비용"],"sections":{"1. 제안 개요":"내용...","2. 현황 분석":"내용...","3. 제안 솔루션":"내용...","4. 수행 방안":"내용...","5. 수행 실적 및 투입 인력":"내용...","6. 유지보수 및 기술 지원":"내용...","7. 투자 비용":"내용..."}}"""

    user_msg = f"""## 회사 정보
{company_info or '(제안사 정보를 기반으로 합리적으로 추정하여 작성하세요)'}

## 레퍼런스
{references or '(유사 프로젝트 경험을 합리적으로 추정하여 작성하세요)'}"""
    if rfp_text:
        user_msg += f"\n\n## RFP 원문\n{rfp_text}"

    result = call_ai(system, user_msg, mock_type="proposal")
    proposal_id = str(uuid.uuid4())[:8]
    db.insert_proposal(proposal_id, rfp_id, result)
    if rfp_id:
        db.update_pipeline_step(rfp_id, "proposal", result)
    record_history(rfp_id, "proposal", result)
    log_activity("제안서 생성", f"Proposal #{proposal_id}")
    await ws_manager.notify("제안서 초안 생성", f"Proposal #{proposal_id}")
    return JSONResponse({"proposal_id": proposal_id, "proposal": result})


# ─── 4. 리뷰 AI ───

@app.post("/api/review-proposal")
async def review_proposal(
    rfp_id: str = Form(None),
    proposal_text: str = Form(""),
):
    rfp_text = ""
    if rfp_id and db.rfp_exists(rfp_id):
        meta = db.get_rfp_meta(rfp_id)
        rfp_text = extract_pdf_text(meta["filepath"])[:4000] if meta and Path(meta["filepath"]).exists() else ""

    system = """당신은 제안서 리뷰 전문가(레드팀)입니다.
제안서를 비판적으로 검토하여 반드시 아래 JSON 형식으로만 반환하세요. 마크다운 코드블록 없이 순수 JSON만 출력하세요:
{"overall_score":78,"grade":"B+","review_items":[{"category":"항목","score":80,"status":"양호/보통/보완필요","comment":"코멘트"}],"missing_requirements":["누락사항"],"logic_issues":["불일치사항"],"improvement_suggestions":["개선제안"]}"""

    user_msg = f"제안서 내용:\n{proposal_text[:6000]}"
    if rfp_text:
        user_msg += f"\n\nRFP 원문:\n{rfp_text}"

    result = call_ai(system, user_msg, mock_type="review")
    if rfp_id:
        db.update_pipeline_step(rfp_id, "review", result)
    record_history(rfp_id, "review", result)
    log_activity("레드팀 리뷰", "제안서 검토 완료")
    return JSONResponse({"review": result})


# ─── 5. Go/No-Go ───

@app.post("/api/strategy")
async def strategy(
    rfp_id: str = Form(None),
    company_strengths: str = Form(""),
    market_context: str = Form(""),
):
    rfp_text = ""
    if rfp_id and db.rfp_exists(rfp_id):
        meta = db.get_rfp_meta(rfp_id)
        rfp_text = extract_pdf_text(meta["filepath"])[:4000] if meta and Path(meta["filepath"]).exists() else ""

    system = """당신은 RFP Go/No-Go 의사결정 전문가입니다.
아래 5개 항목을 반드시 모두 채점하고, 정해진 판정 기준에 따라 GO 또는 NO-GO를 결정하세요.

■ 채점 항목 (각 0~100점):
1. market_fit (시장 적합성): 당사 핵심 역량과 RFP 요구사항의 일치도
   - 80점 이상: 핵심 기술 스택이 일치하고 관련 경험 풍부
   - 60~79점: 부분 일치, 일부 기술 보완 필요
   - 60점 미만: 핵심 역량과 거리가 있음
2. competition (경쟁 환경): 예상 경쟁 강도와 수주 가능성
   - 80점 이상: 경쟁사 적거나 당사 우위 명확
   - 60~79점: 3~5개 경쟁사, 당사 중상위 경쟁력
   - 60점 미만: 강력한 경쟁사 다수, 열위
3. profitability (수익성): 예상 마진율과 비용 대비 수익
   - 80점 이상: 마진 20% 이상 예상
   - 60~79점: 마진 10~20% 예상
   - 60점 미만: 마진 10% 미만 또는 적자 위험
4. strategic_value (전략적 가치): 레퍼런스, 시장 확대, 장기 가치
   - 80점 이상: 핵심 산업군 진출 또는 대형 레퍼런스 확보
   - 60~79점: 일정 수준의 전략적 의미
   - 60점 미만: 전략적 가치 낮음
5. resource_availability (자원 가용성): 투입 인력/기술 확보 가능성
   - 80점 이상: 핵심 인력 즉시 투입 가능
   - 60~79점: 일부 충원/파트너 필요
   - 60점 미만: 핵심 자원 확보 어려움

■ GO/NO-GO 판정 기준 (반드시 준수):
- 5개 항목의 가중 평균 = market_fit×25% + competition×20% + profitability×20% + strategic_value×20% + resource_availability×15%
- 가중 평균 70점 이상 → GO
- 가중 평균 50~69점 → CONDITIONAL GO (조건부 참여)
- 가중 평균 50점 미만 → NO-GO
- confidence(신뢰도)는 입력 정보의 충분성에 따라 산정: 정보 충분→80~95%, 보통→60~79%, 부족→40~59%

■ 중요:
- 동일한 RFP와 입력에 대해서는 반드시 일관된 결과를 내야 함
- 회사 강점이나 시장 상황이 미입력이면 해당 항목을 보수적(50~60점)으로 채점하고, confidence를 낮게(40~60%) 설정
- 감으로 판단하지 말고, 위 기준을 기계적으로 적용하여 채점

반드시 아래 JSON 형식으로만 반환. 마크다운 코드블록 없이 순수 JSON만 출력:
{"recommendation":"GO 또는 CONDITIONAL GO 또는 NO-GO","confidence":"가중평균 기반 신뢰도%","weighted_score":72.5,"analysis":{"market_fit":{"score":80,"weight":"25%","comment":"RFP 원문 근거 기반 설명"},"competition":{"score":65,"weight":"20%","comment":"근거 기반 설명"},"profitability":{"score":70,"weight":"20%","comment":"근거 기반 설명"},"strategic_value":{"score":85,"weight":"20%","comment":"근거 기반 설명"},"resource_availability":{"score":75,"weight":"15%","comment":"근거 기반 설명"}},"key_factors":["RFP 원문에서 도출한 핵심 요인"],"win_strategy":["구체적 수주 전략"],"risks":["참여 시 리스크"]}"""

    user_msg = f"회사 강점: {company_strengths or '(미입력 — 보수적으로 채점)'}\n시장 상황: {market_context or '(미입력 — 보수적으로 채점)'}"
    if rfp_text:
        user_msg += f"\n\nRFP 내용:\n{rfp_text}"

    result = call_ai(system, user_msg, mock_type="strategy")
    if rfp_id:
        db.update_pipeline_step(rfp_id, "strategy", result)
    record_history(rfp_id, "strategy", result)
    log_activity("Go/No-Go 분석", "전략 분석 완료")
    return JSONResponse({"strategy": result})


# ─── 6. 지식 자산화 ───

@app.post("/api/knowledge/save")
async def save_knowledge(
    category: str = Form(...),
    title: str = Form(...),
    content: str = Form(...),
    tags: str = Form(""),
):
    item = {
        "id": str(uuid.uuid4())[:8],
        "category": category,
        "title": title,
        "content": content,
        "tags": [t.strip() for t in tags.split(",") if t.strip()],
        "created_at": datetime.now().isoformat(),
    }
    db.insert_knowledge(item)
    log_activity("지식 저장", title)
    return JSONResponse({"saved": True, "item": item, "total": db.count_knowledge()})


@app.get("/api/knowledge/list")
async def list_knowledge_api(category: Optional[str] = None):
    items = db.list_knowledge(category)
    return JSONResponse({"items": items, "total": len(items)})


@app.post("/api/knowledge/recommend")
async def recommend_knowledge(
    rfp_id: str = Form(None),
    query: str = Form(""),
):
    system = "당신은 지식 자산 관리 전문가입니다. 관련 지식 자산을 추천하고, 재사용 가능한 문장과 템플릿을 JSON으로 제안하세요."
    rfp_text = ""
    if rfp_id and db.rfp_exists(rfp_id):
        meta = db.get_rfp_meta(rfp_id)
        rfp_text = extract_pdf_text(meta["filepath"])[:3000] if meta and Path(meta["filepath"]).exists() else ""
    knowledge_items = db.list_knowledge()
    user_msg = f"검색 쿼리: {query}\n저장된 지식: {json.dumps(knowledge_items, ensure_ascii=False, default=str)[:3000]}"
    if rfp_text:
        user_msg += f"\nRFP: {rfp_text}"
    result = call_ai(system, user_msg, mock_type="knowledge")
    return JSONResponse({"recommendations": result})


# ─── 7. 제안서 Export (Markdown) ───

@app.post("/api/export-proposal")
async def export_proposal(proposal_text: str = Form("")):
    system = """당신은 문서 변환 전문가입니다.
제안서 내용을 깔끔한 Markdown 문서로 변환하세요.
목차, 섹션 제목, 본문을 포함한 완성된 문서를 만들어주세요.
마크다운 코드블록 감싸기 없이 순수 마크다운으로 출력하세요."""

    result = call_ai(system, f"다음 제안서를 Markdown으로 변환:\n\n{proposal_text[:8000]}", mock_type="")
    if not result or result == "분석이 완료되었습니다.":
        result = f"# 제안서\n\n{proposal_text}"
    log_activity("제안서 Export", "Markdown 변환")
    return JSONResponse({"markdown": result})


@app.post("/api/export-docx")
async def export_docx(proposal_text: str = Form("")):
    import io
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    parsed = None
    try:
        stripped = proposal_text.replace("```json", "").replace("```", "").strip()
        parsed = json.loads(stripped)
    except Exception:
        match = re.search(r'\{[\s\S]*\}', stripped)
        if match:
            try:
                parsed = json.loads(match.group(0))
            except Exception:
                pass

    doc = Document()

    # Style defaults
    style = doc.styles['Normal']
    style.font.name = 'Malgun Gothic'
    style.font.size = Pt(11)

    if parsed:
        # Title
        title = parsed.get("title", "제안서")
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(title)
        run.bold = True
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(30, 58, 138)
        doc.add_paragraph()

        # TOC
        toc = parsed.get("table_of_contents", [])
        if toc:
            h = doc.add_heading("목차", level=1)
            h.runs[0].font.color.rgb = RGBColor(30, 58, 138)
            for item in toc:
                level = 1 if item.startswith("  ") else 0
                p = doc.add_paragraph(item.strip(), style='List Bullet' if level else 'List Number')
            doc.add_page_break()

        # Sections
        sections = parsed.get("sections", {})
        for sec_title, sec_content in sections.items():
            h = doc.add_heading(sec_title, level=2)
            h.runs[0].font.color.rgb = RGBColor(30, 58, 138)
            doc.add_paragraph(sec_content)
            doc.add_paragraph()
    else:
        # Plain text fallback
        doc.add_heading("제안서", level=1)
        for line in proposal_text.split("\n"):
            line = line.strip()
            if not line:
                doc.add_paragraph()
            elif line.startswith("#"):
                level = min(line.count("#"), 4)
                doc.add_heading(line.lstrip("# "), level=level)
            else:
                doc.add_paragraph(line)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    log_activity("제안서 Export", "DOCX 다운로드")
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": "attachment; filename=proposal.docx"},
    )


@app.post("/api/export-pptx")
async def export_pptx(proposal_text: str = Form("")):
    import io, math
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as PptRGB
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    parsed = None
    try:
        stripped = proposal_text.replace("```json", "").replace("```", "").strip()
        parsed = json.loads(stripped)
    except Exception:
        match = re.search(r'\{[\s\S]*\}', stripped)
        if match:
            try: parsed = json.loads(match.group(0))
            except Exception: pass

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    WHITE = PptRGB(255, 255, 255)
    GRAY = PptRGB(148, 163, 184)
    DARK_BG = PptRGB(15, 23, 42)
    CARD_BG = PptRGB(30, 41, 59)
    ACCENT = PptRGB(6, 182, 212)
    LIGHT_BLUE = PptRGB(99, 102, 241)
    GREEN = PptRGB(16, 185, 129)
    AMBER = PptRGB(245, 158, 11)
    PURPLE = PptRGB(139, 92, 246)
    CARD_COLORS = [LIGHT_BLUE, ACCENT, GREEN, AMBER, PURPLE, PptRGB(236, 72, 153)]

    def add_bg(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = DARK_BG

    def rect(slide, l, t, w, h, color, radius=False):
        shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        s = slide.shapes.add_shape(shp_type, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
        return s

    def oval(slide, l, t, w, h, color):
        s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
        return s

    def arrow(slide, l, t, w, h, color):
        s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
        return s

    def chevron(slide, l, t, w, h, color):
        s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
        return s

    def tbox(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
        p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
        return tb

    def shaped_text(shape, text, sz=12, color=WHITE, bold=False, align=PP_ALIGN.CENTER):
        tf = shape.text_frame; tf.word_wrap = True
        tf.paragraphs[0].alignment = align
        p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
        p.font.color.rgb = color; p.font.bold = bold

    def slide_header(slide, title):
        rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), LIGHT_BLUE)
        tbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.8), title, sz=28, bold=True)
        rect(slide, Inches(0.7), Inches(1.15), Inches(1.5), Inches(0.05), ACCENT)

    # Max chars per slide content area (~13pt, 11 inch wide) ≈ 380 chars
    MAX_CHARS = 350

    def split_text(text, limit=MAX_CHARS):
        """Split text into chunks that fit one slide."""
        sentences = re.split(r'(?<=[.!?。])\s+', text)
        chunks, cur = [], ""
        for s in sentences:
            if len(cur) + len(s) + 1 > limit and cur:
                chunks.append(cur.strip())
                cur = s
            else:
                cur = cur + " " + s if cur else s
        if cur.strip(): chunks.append(cur.strip())
        return chunks if chunks else [text[:limit]]

    # Detect if content is suitable for diagram
    def detect_visual_type(title, content):
        tl = (title + " " + content).lower()
        if any(k in tl for k in ["아키텍처", "architecture", "시스템 구성", "구조", "플랫폼"]):
            return "architecture"
        if any(k in tl for k in ["일정", "schedule", "단계", "phase", "로드맵", "마일스톤", "timeline"]):
            return "timeline"
        if any(k in tl for k in ["프로세스", "절차", "workflow", "흐름", "방법론", "접근"]):
            return "process"
        if any(k in tl for k in ["장점", "강점", "특징", "핵심", "차별", "benefit", "advantage"]):
            return "cards"
        if any(k in tl for k in ["조직", "인력", "체계", "팀", "역할"]):
            return "org"
        return "text"

    def extract_items(content):
        """Extract bullet-like items from text."""
        items = []
        for line in content.replace(".", ".\n").split("\n"):
            line = line.strip().lstrip("-•·▶▷◆ ")
            if len(line) > 5: items.append(line)
        return items[:8]  # max 8 items for visual

    def add_architecture_slide(slide, title, content):
        slide_header(slide, title)
        items = extract_items(content)
        if len(items) < 3: items = [content[i:i+40] for i in range(0, min(len(content), 200), 40)]
        # Draw layered architecture
        layers = items[:5]
        y_start = Inches(1.6)
        for i, layer in enumerate(layers):
            c = CARD_COLORS[i % len(CARD_COLORS)]
            bg_c = PptRGB(c.red // 4 + 10, c.green // 4 + 10, c.blue // 4 + 10) if hasattr(c, 'red') else CARD_BG
            s = rect(slide, Inches(1.5), y_start + Inches(i * 1.05), Inches(10), Inches(0.9), CARD_BG, radius=True)
            # Color accent bar on left
            rect(slide, Inches(1.5), y_start + Inches(i * 1.05), Inches(0.12), Inches(0.9), c)
            tbox(slide, Inches(1.9), y_start + Inches(i * 1.05) + Inches(0.15), Inches(9), Inches(0.6), layer[:80], sz=14, color=WHITE, bold=True if i == 0 else False)
        # Side label
        tbox(slide, Inches(0.7), Inches(6.5), Inches(5), Inches(0.4), f"{len(layers)}개 계층 아키텍처", sz=11, color=GRAY)

    def add_timeline_slide(slide, title, content):
        slide_header(slide, title)
        items = extract_items(content)
        if len(items) < 2: items = content.split(",")
        items = [it.strip() for it in items if it.strip()][:6]
        n = len(items)
        # Horizontal timeline
        y_line = Inches(3.5)
        rect(slide, Inches(1), y_line, Inches(11), Inches(0.06), PptRGB(51, 65, 85))
        for i, item in enumerate(n and items or ["Phase 1"]):
            x = Inches(1.2 + i * (10.5 / max(n - 1, 1))) if n > 1 else Inches(6)
            c = CARD_COLORS[i % len(CARD_COLORS)]
            # Circle node
            o = oval(slide, x - Inches(0.2), y_line - Inches(0.17), Inches(0.4), Inches(0.4), c)
            shaped_text(o, str(i + 1), sz=11, bold=True)
            # Label above
            tbox(slide, x - Inches(1.2), y_line - Inches(1.3), Inches(2.4), Inches(1.0), item[:50], sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            # Phase label below
            tbox(slide, x - Inches(0.8), y_line + Inches(0.5), Inches(1.6), Inches(0.4), f"Phase {i + 1}", sz=10, color=GRAY, align=PP_ALIGN.CENTER)

    def add_process_slide(slide, title, content):
        slide_header(slide, title)
        items = extract_items(content)[:5]
        if len(items) < 2: items = ["분석", "설계", "구현", "테스트", "배포"]
        n = len(items)
        # Chevron process flow
        chev_w = Inches(min(2.2, 11 / n))
        gap = Inches(0.1)
        total_w = n * (chev_w + gap)
        start_x = (Inches(13.333) - total_w) // 2 + Inches(0.3)
        for i, item in enumerate(items):
            x = start_x + i * (chev_w + gap)
            c = CARD_COLORS[i % len(CARD_COLORS)]
            s = chevron(slide, x, Inches(2.8), chev_w, Inches(1.6), c)
            shaped_text(s, item[:20], sz=13, bold=True)
            # Description below
            tbox(slide, x, Inches(4.6), chev_w, Inches(1.2), "", sz=10, color=GRAY, align=PP_ALIGN.CENTER)

    def add_cards_slide(slide, title, content):
        slide_header(slide, title)
        items = extract_items(content)[:6]
        if len(items) < 2: items = [content[:60]]
        n = len(items)
        cols = min(n, 3)
        rows = math.ceil(n / cols)
        card_w = Inches(3.4)
        card_h = Inches(2.0)
        gap_x = Inches(0.4)
        gap_y = Inches(0.3)
        total_w = cols * card_w + (cols - 1) * gap_x
        start_x = (Inches(13.333) - total_w) // 2
        for i, item in enumerate(items):
            col = i % cols
            row = i // cols
            x = start_x + col * (card_w + gap_x)
            y = Inches(1.6) + row * (card_h + gap_y)
            c = CARD_COLORS[i % len(CARD_COLORS)]
            # Card bg
            card = rect(slide, x, y, card_w, card_h, CARD_BG, radius=True)
            # Top accent
            rect(slide, x, y, card_w, Inches(0.06), c)
            # Number circle
            num_s = oval(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.45), Inches(0.45), c)
            shaped_text(num_s, str(i + 1), sz=14, bold=True)
            # Text
            tbox(slide, x + Inches(0.85), y + Inches(0.3), card_w - Inches(1.1), card_h - Inches(0.5), item[:80], sz=13, color=WHITE)

    def add_org_slide(slide, title, content):
        slide_header(slide, title)
        items = extract_items(content)[:6]
        if not items: items = ["프로젝트 관리"]
        # Top box (PM)
        pm_w = Inches(3)
        pm_x = (Inches(13.333) - pm_w) // 2
        s = rect(slide, pm_x, Inches(1.6), pm_w, Inches(0.9), LIGHT_BLUE, radius=True)
        shaped_text(s, items[0][:30], sz=14, bold=True)
        # Lines + sub boxes
        subs = items[1:] if len(items) > 1 else ["기술팀", "기획팀", "QA팀"]
        n = len(subs)
        sub_w = Inches(2.4)
        gap = Inches(0.3)
        total = n * sub_w + (n - 1) * gap
        start_x = (Inches(13.333) - total) // 2
        center_x = pm_x + pm_w // 2
        rect(slide, center_x - Inches(0.02), Inches(2.5), Inches(0.04), Inches(0.6), PptRGB(51, 65, 85))
        for i, sub in enumerate(subs):
            x = start_x + i * (sub_w + gap)
            c = CARD_COLORS[(i + 1) % len(CARD_COLORS)]
            s = rect(slide, x, Inches(3.3), sub_w, Inches(0.8), c, radius=True)
            shaped_text(s, sub[:25], sz=13, bold=True)
            # Connect line
            cx = x + sub_w // 2
            rect(slide, cx - Inches(0.02), Inches(3.1), Inches(0.04), Inches(0.2), PptRGB(51, 65, 85))

    def add_text_slide(slide, title, text_chunk):
        slide_header(slide, title)
        card = rect(slide, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5), CARD_BG, radius=True)
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.5), Inches(4.9))
        tf = tb.text_frame; tf.word_wrap = True
        for j, line in enumerate(text_chunk.split("\n")):
            line = line.strip()
            if not line: continue
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = line[:120]  # hard limit per line
            p.font.size = Pt(14); p.font.color.rgb = PptRGB(203, 213, 225)
            p.space_after = Pt(8); p.line_spacing = Pt(22)

    # ─── Build slides ───
    title_text = parsed.get("title", "제안서") if parsed else "제안서"
    toc_items = parsed.get("table_of_contents", []) if parsed else []
    sections = parsed.get("sections", {}) if parsed else {}

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), LIGHT_BLUE)
    oval(slide, Inches(9.5), Inches(-0.5), Inches(4.5), Inches(4.5), PptRGB(30, 41, 59))
    oval(slide, Inches(10), Inches(4.5), Inches(3), Inches(3), PptRGB(25, 35, 52))
    tbox(slide, Inches(1), Inches(2.2), Inches(9), Inches(1.5), title_text, sz=40, bold=True)
    rect(slide, Inches(1), Inches(3.8), Inches(2), Inches(0.06), ACCENT)
    from datetime import datetime as dt
    tbox(slide, Inches(1), Inches(4.1), Inches(8), Inches(0.8), "AI-Powered Proposal | PARAIX Hackathon 2026", sz=16, color=GRAY)
    tbox(slide, Inches(1), Inches(5.5), Inches(5), Inches(0.5), dt.now().strftime("%Y년 %m월 %d일"), sz=14, color=GRAY)

    # Slide 2: TOC
    if toc_items:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide); slide_header(slide, "목차")
        y = Inches(1.6)
        for i, item in enumerate(toc_items):
            is_sub = item.startswith("  ")
            left = Inches(1.4) if is_sub else Inches(0.7)
            sz = 13 if is_sub else 17
            clr = GRAY if is_sub else WHITE
            tbox(slide, left, y, Inches(10), Inches(0.42), item.strip(), sz=sz, color=clr, bold=not is_sub)
            y += Inches(0.35) if is_sub else Inches(0.42)
            if y > Inches(6.6):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_bg(slide); slide_header(slide, "목차 (계속)")
                y = Inches(1.6)

    # Section slides with auto visual detection
    for sec_title, sec_content in sections.items():
        content = sec_content if isinstance(sec_content, str) else str(sec_content)
        vis_type = detect_visual_type(sec_title, content)

        if vis_type == "architecture":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide); add_architecture_slide(slide, sec_title, content)
        elif vis_type == "timeline":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide); add_timeline_slide(slide, sec_title, content)
        elif vis_type == "process":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide); add_process_slide(slide, sec_title, content)
        elif vis_type == "cards":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide); add_cards_slide(slide, sec_title, content)
        elif vis_type == "org":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide); add_org_slide(slide, sec_title, content)
        else:
            # Text with auto-split across slides
            chunks = split_text(content)
            for ci, chunk in enumerate(chunks):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_bg(slide)
                label = sec_title if ci == 0 else f"{sec_title} (계속)"
                add_text_slide(slide, label, chunk)

    # Last slide: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    oval(slide, Inches(4), Inches(0.5), Inches(5.333), Inches(5.333), PptRGB(20, 30, 48))
    tbox(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(1.2), "Thank You", sz=48, bold=True, align=PP_ALIGN.CENTER)
    tbox(slide, Inches(0), Inches(3.8), Inches(13.333), Inches(0.6), "AI-Powered by RFP AI Analyzer", sz=18, color=GRAY, align=PP_ALIGN.CENTER)
    rect(slide, Inches(5.5), Inches(4.5), Inches(2.333), Inches(0.05), ACCENT)

    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    log_activity("제안서 Export", "PPTX 다운로드")
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=proposal.pptx"})


# ─── 8. 사업 견적 산출 ───

@app.post("/api/estimate")
async def estimate_cost(rfp_id: str = Form(None), additional_info: str = Form("")):
    rfp_text = ""
    if rfp_id and db.rfp_exists(rfp_id):
        meta = db.get_rfp_meta(rfp_id)
        rfp_text = extract_pdf_text(meta["filepath"])[:6000] if meta and Path(meta["filepath"]).exists() else ""

    system = """당신은 IT 프로젝트 사업 견적 전문가입니다.
RFP를 분석하여 상세한 사업 견적을 산출하세요.

반드시 아래 JSON 형식으로만 반환하세요. 마크다운 코드블록 없이 순수 JSON만 출력:
{
  "project_name": "프로젝트명",
  "total_cost": "총 예상 금액 (원)",
  "total_cost_number": 0,
  "duration_months": 6,
  "summary": "견적 요약 설명 (2~3문장)",
  "categories": [
    {
      "name": "인건비",
      "subtotal": "금액 (원)",
      "subtotal_number": 0,
      "ratio": "전체 비율 (%)",
      "items": [
        {"role": "역할명", "grade": "등급 (특급/고급/중급/초급)", "count": 1, "months": 6, "unit_cost": "월 단가 (원)", "cost": "금액 (원)", "reason": "해당 인력이 필요한 상세 사유"}
      ]
    },
    {
      "name": "SW 라이선스",
      "subtotal": "금액 (원)",
      "subtotal_number": 0,
      "ratio": "비율 (%)",
      "items": [
        {"item": "항목명", "cost": "금액 (원)", "reason": "필요 사유"}
      ]
    },
    {
      "name": "HW/인프라",
      "subtotal": "금액 (원)",
      "subtotal_number": 0,
      "ratio": "비율 (%)",
      "items": [
        {"item": "항목명", "cost": "금액 (원)", "reason": "필요 사유"}
      ]
    },
    {
      "name": "기타 경비",
      "subtotal": "금액 (원)",
      "subtotal_number": 0,
      "ratio": "비율 (%)",
      "items": [
        {"item": "항목명", "cost": "금액 (원)", "reason": "필요 사유"}
      ]
    }
  ],
  "risks": [
    {"risk": "비용 리스크", "impact": "영향 금액 또는 비율", "mitigation": "대응 방안"}
  ],
  "assumptions": ["견적 전제 조건 1", "견적 전제 조건 2"],
  "notes": "참고사항 및 할인/협상 여지"
}

산출 기준 (반드시 준수):
- 인건비: SW기술자 노임단가 기준 (2024년 기준 특급 약 450만원/월, 고급 380만원, 중급 310만원, 초급 240만원)
- RFP 요구사항의 복잡도, 기술 스택, 기간을 고려하여 현실적인 인력 구성

HW/인프라 가격 기준 (반드시 준수):
- 클라우드 서비스(AWS/Azure/GCP)를 우선 검토. 온프레미스 서버 구매는 RFP가 명시적으로 요구할 때만
- GPU 서버: 클라우드 GPU 인스턴스 임대 기준으로 산정 (예: AWS p4d.24xlarge 약 월 3,000~4,000만원)
- 온프레미스 GPU 서버 구매가 필요한 경우: NVIDIA H100 서버 1대 약 3~5억원, A100 서버 1대 약 1~2억원 수준
- 일반 서버: 클라우드 기준 월 50~300만원, 온프레미스 구매 시 대당 1,000~5,000만원
- 네트워크/보안 장비: 방화벽 1,000~3,000만원, L4 스위치 500~1,500만원

금액 검증 규칙:
- 단일 항목이 전체 견적의 60%를 초과하면 안 됨
- HW/인프라 비용이 전체의 40%를 초과하면 클라우드 대안을 반드시 제시
- 각 항목의 금액은 한국 IT 시장의 실제 거래 가격 기준
- 비현실적으로 높은 금액(단일 항목 10억 이상)은 근거를 반드시 명시
- 총액은 한국 IT 시장 기준 현실적인 수준으로 산정"""

    user_msg = "다음 RFP를 분석하여 상세 사업 견적을 산출해주세요."
    if rfp_text:
        user_msg += f"\n\nRFP 내용:\n{rfp_text}"
    if additional_info:
        user_msg += f"\n\n추가 정보:\n{additional_info}"

    result = call_ai(system, user_msg, mock_type="estimate")
    if rfp_id:
        db.update_pipeline_step(rfp_id, "estimate", result)
    record_history(rfp_id or "", "estimate", result)
    meta = db.get_rfp_meta(rfp_id) if rfp_id else None
    log_activity("견적 산출", meta["filename"] if meta else "")
    return JSONResponse({"estimate": result})


# ─── 9. 원클릭 파이프라인 ───

@app.post("/api/pipeline/run")
async def run_pipeline(
    rfp_id: str = Form(...),
    company_info: str = Form(""),
    industry: str = Form("IT"),
    customer_type: str = Form("대기업"),
):
    rfp = db.get_rfp_meta(rfp_id)
    if not rfp:
        raise HTTPException(404, "RFP를 찾을 수 없습니다.")
    rfp_text = extract_pdf_text(rfp["filepath"])[:6000] if rfp.get("filepath") and Path(rfp["filepath"]).exists() else ""
    results = {}
    steps_done = []

    # Step 1: Analyze
    system1 = """RFP를 분석하여 반드시 JSON으로만 반환하세요. 마크다운 코드블록 없이 순수 JSON만:
{"summary":"요약","requirements":[{"id":"REQ-001","category":"분류","description":"설명","priority":"필수","risk":"높음"}],"evaluation_criteria":[{"criteria":"기준","weight":"30%","description":"설명"}],"risks":[{"risk":"리스크","level":"높음","description":"설명"}]}"""
    results["analyze"] = call_ai(system1, rfp_text, mock_type="analyze")
    steps_done.append("analyze")

    # Step 2: Pattern
    system2 = """Winning Proposal 패턴을 JSON으로만 반환. 마크다운 코드블록 없이:
{"industry_analysis":"분석","winning_patterns":[{"pattern":"패턴","description":"설명","confidence":"높음"}],"style_recommendations":["추천"],"differentiation_tips":["전략"]}"""
    results["pattern"] = call_ai(system2, f"산업:{industry} 고객:{customer_type}\nRFP:{rfp_text[:3000]}", mock_type="pattern")
    steps_done.append("pattern")

    # Step 3: Proposal
    system3 = """제출 가능한 수준의 제안서 초안을 작성하세요. 각 섹션 200자 이상, 구체적 수치 포함, 한국어 비즈니스 문체.
필수 섹션: 제안개요, 현황분석, 제안솔루션, 수행방안, 수행실적, 유지보수, 투자비용.
JSON으로만 반환. 마크다운 코드블록 없이:
{"title":"제목","table_of_contents":["1. 제안 개요","2. 현황 분석","3. 제안 솔루션","4. 수행 방안","5. 수행 실적","6. 유지보수","7. 투자 비용"],"sections":{"1. 제안 개요":"내용","2. 현황 분석":"내용","3. 제안 솔루션":"내용","4. 수행 방안":"내용","5. 수행 실적":"내용","6. 유지보수":"내용","7. 투자 비용":"내용"}}"""
    results["proposal"] = call_ai(system3, f"회사:{company_info}\nRFP:{rfp_text}", mock_type="proposal")
    steps_done.append("proposal")

    # Step 4: Review
    system4 = """제안서를 리뷰하여 JSON으로만 반환. 마크다운 코드블록 없이:
{"overall_score":78,"grade":"B+","review_items":[{"category":"항목","score":80,"status":"양호","comment":"코멘트"}],"missing_requirements":["누락"],"logic_issues":["불일치"],"improvement_suggestions":["개선"]}"""
    results["review"] = call_ai(system4, f"제안서:{results['proposal'][:4000]}\nRFP:{rfp_text[:3000]}", mock_type="review")
    steps_done.append("review")

    # Step 5: Strategy
    system5 = """Go/No-Go를 JSON으로만 반환. 마크다운 코드블록 없이:
{"recommendation":"GO","confidence":"75%","analysis":{"market_fit":{"score":80,"comment":"설명"}},"key_factors":["요인"],"win_strategy":["전략"]}"""
    results["strategy"] = call_ai(system5, f"회사:{company_info}\nRFP:{rfp_text[:3000]}", mock_type="strategy")
    steps_done.append("strategy")

    db.upsert_pipeline(rfp_id, {s: True for s in steps_done}, results)
    for step in steps_done:
        record_history(rfp_id, step, results[step])
    log_activity("파이프라인 완료", f"{rfp['filename']} - {len(steps_done)}단계")
    await ws_manager.notify("파이프라인 완료", f"{rfp['filename']} - {len(steps_done)}단계")

    return JSONResponse({"rfp_id": rfp_id, "steps_completed": steps_done, "results": results})


# ─── 9. AI 경쟁력 채점 ───

@app.post("/api/score-proposal")
async def score_proposal(proposal_text: str = Form("")):
    system = """당신은 제안서 경쟁력 채점 AI입니다.
제안서를 분석하고 반드시 아래 JSON 형식으로만 반환하세요. 마크다운 코드블록 없이 순수 JSON만:
{"total_score":82,"max_score":100,"grade":"A","categories":[{"name":"기술 이해도","score":85,"max":100,"feedback":"피드백"},{"name":"실현 가능성","score":78,"max":100,"feedback":"피드백"},{"name":"차별화","score":70,"max":100,"feedback":"피드백"},{"name":"문서 완성도","score":80,"max":100,"feedback":"피드백"},{"name":"가격 경쟁력","score":75,"max":100,"feedback":"피드백"}],"strengths":["강점1"],"weaknesses":["약점1"],"win_probability":"65%"}"""

    result = call_ai(system, f"다음 제안서의 경쟁력을 채점해주세요:\n\n{proposal_text[:6000]}", mock_type="")
    if not result or "분석이 완료" in result:
        result = json.dumps({
            "total_score": 82, "max_score": 100, "grade": "A-",
            "categories": [
                {"name": "기술 이해도", "score": 85, "max": 100, "feedback": "요구사항에 대한 기술적 이해가 높음"},
                {"name": "실현 가능성", "score": 78, "max": 100, "feedback": "일정 계획이 다소 낙관적"},
                {"name": "차별화", "score": 70, "max": 100, "feedback": "경쟁사 대비 뚜렷한 차별점 보강 필요"},
                {"name": "문서 완성도", "score": 88, "max": 100, "feedback": "체계적 구성, 시각자료 추가 권장"},
                {"name": "가격 경쟁력", "score": 75, "max": 100, "feedback": "시장 평균 수준, 가치 설명 보강 필요"},
            ],
            "strengths": ["기술 역량 충분", "체계적 문서 구성", "풍부한 레퍼런스"],
            "weaknesses": ["차별화 포인트 부족", "일정 리스크", "가격 근거 불충분"],
            "win_probability": "65%",
        }, ensure_ascii=False, indent=2)
    log_activity("경쟁력 채점", "AI 채점 완료")
    return JSONResponse({"score": result})


# ─── 결과 조회 API ───

@app.get("/api/results/{rfp_id}")
async def get_results(rfp_id: str):
    data = db.get_pipeline(rfp_id)
    return JSONResponse({
        "results": data.get("results", {}),
        "steps": list(data.get("completed_steps", {}).keys()),
    })


# ─── 10. 제안서 버전 관리 ───

@app.post("/api/version/save")
async def save_version(
    rfp_id: str = Form(...),
    content: str = Form(...),
    score: int = Form(0),
    note: str = Form(""),
):
    ver_num = db.count_versions(rfp_id) + 1
    created_at = datetime.now().strftime("%Y-%m-%d %H:%M")
    db.insert_version(rfp_id, ver_num, content, score, note, created_at)
    version = {"version": ver_num, "content": content, "score": score, "note": note, "created_at": created_at}
    log_activity("버전 저장", f"v{ver_num} (점수:{score})")
    return JSONResponse({"saved": True, "version": version, "total": ver_num})


@app.get("/api/version/list/{rfp_id}")
async def list_versions_api(rfp_id: str):
    versions = db.list_versions(rfp_id)
    return JSONResponse({"versions": versions})


# ─── 11. 팀 협업 ───

@app.post("/api/team/init")
async def team_init(rfp_id: str = Form(...)):
    return JSONResponse(db.get_team(rfp_id))


@app.post("/api/team/add-member")
async def add_member(rfp_id: str = Form(...), name: str = Form(...), role: str = Form("")):
    member_id = str(uuid.uuid4())[:6]
    db.add_team_member(rfp_id, member_id, name, role)
    log_activity("팀원 추가", f"{name} ({role})")
    await ws_manager.notify("팀원 추가", f"{name} ({role})")
    return JSONResponse({"members": db.get_team_members(rfp_id)})


@app.post("/api/team/remove-member")
async def remove_member(rfp_id: str = Form(...), member_id: str = Form(...)):
    existed = db.remove_team_member(rfp_id, member_id)
    if existed:
        log_activity("팀원 삭제", f"ID: {member_id}")
        await ws_manager.notify("팀원 삭제", member_id)
    return JSONResponse({"members": db.get_team_members(rfp_id)})


@app.post("/api/team/add-section")
async def add_section(
    rfp_id: str = Form(...),
    title: str = Form(...),
    assignee: str = Form(""),
):
    section_id = str(uuid.uuid4())[:6]
    db.add_team_section(rfp_id, section_id, title, assignee)
    team = db.get_team(rfp_id)
    return JSONResponse({"sections": team["sections"]})


@app.post("/api/team/update-assignee")
async def update_assignee(
    rfp_id: str = Form(...),
    section_id: str = Form(...),
    assignee: str = Form(...),
):
    db.update_section_assignee(section_id, assignee)
    return JSONResponse({"ok": True})


@app.post("/api/team/auto-assign")
async def auto_assign(
    rfp_id: str = Form(...),
):
    if not db.rfp_exists(rfp_id):
        raise HTTPException(404, "RFP를 찾을 수 없습니다.")
    meta = db.get_rfp_meta(rfp_id)
    rfp_text = extract_pdf_text(meta["filepath"])[:6000] if meta and Path(meta["filepath"]).exists() else ""
    members = db.get_team_members(rfp_id)
    member_info = ", ".join([f"{m['name']}({m['role']})" for m in members]) if members else "(팀원 미등록 — 역할명으로 배정)"

    system = f"""당신은 프로젝트 매니저입니다.
RFP를 분석하여 제안서 작성에 필요한 섹션과 담당자를 자동 배정하세요.

등록된 팀원: {member_info}

규칙:
1. RFP 요구사항을 기반으로 제안서에 필요한 섹션(7~12개)을 도출
2. 각 섹션에 가장 적합한 팀원(또는 역할)을 배정
3. 팀원이 없으면 "PM", "기술리드", "아키텍트", "기획자", "디자이너" 등 역할명으로 배정

반드시 아래 JSON 형식으로만 반환. 마크다운 코드블록 없이 순수 JSON:
{{"sections":[{{"title":"섹션명","assignee":"담당자명","reason":"배정 사유"}}]}}"""

    result = call_ai(system, rfp_text, mock_type="")
    if not result or "분석이 완료" in result:
        result = json.dumps({"sections": [
            {"title": "1. 제안 개요", "assignee": "PM", "reason": "전체 방향 설정"},
            {"title": "2. 현황 분석", "assignee": "기획자", "reason": "고객 환경 분석"},
            {"title": "3. 시스템 아키텍처", "assignee": "아키텍트", "reason": "기술 설계"},
            {"title": "4. 핵심 기능 상세", "assignee": "기술리드", "reason": "기능 구현 방안"},
            {"title": "5. 수행 방안", "assignee": "PM", "reason": "일정/품질 관리"},
            {"title": "6. 수행 실적", "assignee": "기획자", "reason": "레퍼런스 정리"},
            {"title": "7. 투입 인력", "assignee": "PM", "reason": "인력 계획"},
            {"title": "8. 유지보수", "assignee": "기술리드", "reason": "기술 지원 계획"},
            {"title": "9. 투자 비용", "assignee": "PM", "reason": "비용 산정"},
        ]}, ensure_ascii=False)

    # Parse and apply
    parsed = None
    try:
        stripped = result.replace("```json", "").replace("```", "").strip()
        parsed = json.loads(stripped)
    except Exception:
        match = re.search(r'\{[\s\S]*\}', result)
        if match:
            try:
                parsed = json.loads(match.group(0))
            except Exception:
                pass

    new_sections = []
    if parsed and "sections" in parsed:
        for s in parsed["sections"]:
            new_sections.append({
                "id": str(uuid.uuid4())[:6],
                "title": s.get("title", ""),
                "assignee": s.get("assignee", ""),
                "reason": s.get("reason", ""),
                "status": "대기",
                "comments": [],
            })
        db.replace_team_sections(rfp_id, new_sections)

    log_activity("AI 자동 배정", f"{len(new_sections)}개 섹션")
    return JSONResponse({"sections": new_sections, "raw": result})


@app.post("/api/team/update-status")
async def update_section_status(
    rfp_id: str = Form(...),
    section_id: str = Form(...),
    status: str = Form(...),
):
    title = db.update_section_status(section_id, status)
    if title:
        await ws_manager.notify("섹션 상태 변경", f"{title} → {status}")
    return JSONResponse({"ok": True})


@app.post("/api/team/add-comment")
async def add_comment(
    rfp_id: str = Form(...),
    section_id: str = Form(...),
    author: str = Form(...),
    text: str = Form(...),
):
    db.add_section_comment(section_id, author, text, datetime.now().strftime("%H:%M"))
    await ws_manager.notify("코멘트 추가", f"{author}: {text[:30]}")
    return JSONResponse({"ok": True})


@app.get("/api/team/{rfp_id}")
async def get_team_api(rfp_id: str):
    return JSONResponse(db.get_team(rfp_id))


# ─── 12. 자동 일정 생성 ───

@app.post("/api/schedule/generate")
async def generate_schedule(
    deadline: str = Form(...),
    rfp_id: str = Form(None),
):
    try:
        dl = datetime.strptime(deadline, "%Y-%m-%d")
    except ValueError:
        raise HTTPException(400, "날짜 형식: YYYY-MM-DD")

    today = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)
    total_days = (dl - today).days
    if total_days < 1:
        raise HTTPException(400, "마감일은 오늘 이후여야 합니다.")

    phases = [
        {"name": "RFP 분석 및 전략 수립", "ratio": 0.15, "tasks": ["RFP 업로드/구조화 분석", "Go/No-Go 의사결정", "Winning 패턴 분석", "팀 구성 및 역할 배정"]},
        {"name": "제안서 초안 작성", "ratio": 0.35, "tasks": ["AI 초안 생성", "섹션별 담당자 작성", "기술 솔루션 상세화", "레퍼런스/사례 정리"]},
        {"name": "내부 리뷰 및 수정", "ratio": 0.25, "tasks": ["Red Team AI 리뷰", "요구사항 누락 체크", "논리 일관성 검토", "경쟁력 채점 및 보완"]},
        {"name": "최종 검수 및 제출", "ratio": 0.15, "tasks": ["디자인/레이아웃 최종화", "PDF/DOCX 생성", "경영진 승인", "제출"]},
        {"name": "버퍼 (예비)", "ratio": 0.10, "tasks": ["긴급 수정 대응", "최종 검토"]},
    ]

    schedule = []
    current = today
    for phase in phases:
        days = max(1, round(total_days * phase["ratio"]))
        end = min(current + timedelta(days=days), dl)
        schedule.append({
            "phase": phase["name"],
            "start": current.strftime("%Y-%m-%d"),
            "end": end.strftime("%Y-%m-%d"),
            "days": (end - current).days,
            "tasks": phase["tasks"],
        })
        current = end

    log_activity("일정 생성", f"마감: {deadline} ({total_days}일)")
    return JSONResponse({
        "deadline": deadline,
        "total_days": total_days,
        "schedule": schedule,
    })


# ─── 13. PDF 제안서 생성 ───

@app.post("/api/export-pdf")
async def export_pdf(proposal_text: str = Form("")):
    import io
    from fpdf import FPDF

    parsed = None
    try:
        stripped = proposal_text.replace("```json", "").replace("```", "").strip()
        parsed = json.loads(stripped)
    except Exception:
        match = re.search(r'\{[\s\S]*\}', proposal_text)
        if match:
            try:
                parsed = json.loads(match.group(0))
            except Exception:
                pass

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)

    # Load Korean font
    font_path = Path("/usr/share/fonts/truetype/nanum/NanumGothic.ttf")
    has_korean_font = font_path.exists()
    if has_korean_font:
        pdf.add_font("Nanum", "", str(font_path), uni=True)
        pdf.add_font("Nanum", "B", str(font_path.parent / "NanumGothicBold.ttf"), uni=True)
        body_font, bold_font = "Nanum", "Nanum"
    else:
        body_font, bold_font = "Helvetica", "Helvetica"

    # Cover page
    pdf.add_page()
    pdf.set_fill_color(30, 58, 138)
    pdf.rect(0, 0, 210, 100, 'F')
    pdf.set_y(30)
    pdf.set_font(bold_font, "B", 28)
    pdf.set_text_color(255, 255, 255)
    title = parsed.get("title", "Proposal") if parsed else "Proposal"
    pdf.cell(0, 15, title, ln=True, align="C")
    pdf.set_font(body_font, "", 14)
    pdf.cell(0, 10, datetime.now().strftime("%Y-%m-%d"), ln=True, align="C")
    pdf.set_y(110)
    pdf.set_text_color(60, 60, 60)
    pdf.set_font(body_font, "", 11)
    pdf.cell(0, 8, "Generated by RFP AI Analyzer", ln=True, align="C")

    if parsed:
        # TOC page
        toc = parsed.get("table_of_contents", [])
        if toc:
            pdf.add_page()
            pdf.set_font(bold_font, "B", 18)
            pdf.set_text_color(30, 58, 138)
            pdf.cell(0, 12, "Table of Contents", ln=True)
            pdf.ln(4)
            pdf.set_font(body_font, "", 12)
            pdf.set_text_color(50, 50, 50)
            for item in toc:
                prefix = "    " if item.startswith("  ") else ""
                pdf.cell(0, 8, f"{prefix}{item.strip()}", ln=True)

        # Sections
        sections = parsed.get("sections", {})
        for sec_title, sec_content in sections.items():
            pdf.add_page()
            pdf.set_font(bold_font, "B", 16)
            pdf.set_text_color(30, 58, 138)
            pdf.cell(0, 12, sec_title, ln=True)
            pdf.ln(2)
            pdf.set_draw_color(30, 58, 138)
            pdf.line(10, pdf.get_y(), 200, pdf.get_y())
            pdf.ln(6)
            pdf.set_font(body_font, "", 11)
            pdf.set_text_color(50, 50, 50)
            pdf.multi_cell(0, 7, sec_content)
    else:
        pdf.add_page()
        pdf.set_font(body_font, "", 11)
        pdf.set_text_color(50, 50, 50)
        pdf.multi_cell(0, 7, proposal_text[:5000])

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)

    log_activity("PDF Export", "제안서 PDF 생성")
    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": "attachment; filename=proposal.pdf"},
    )


@app.get("/api/history/{rfp_id}")
async def get_history_api(rfp_id: str, step: str = None):
    items = db.get_history(rfp_id, step)
    return JSONResponse({"history": items})


@app.get("/api/rfp-list")
async def rfp_list_api():
    rfps = db.list_rfps()
    pipelines = db.list_pipelines()
    items = []
    for v in rfps:
        steps = list(pipelines.get(v["id"], {}).get("completed_steps", {}).keys())
        items.append({"id": v["id"], "filename": v["filename"], "text_length": v["text_length"],
                       "uploaded_at": v.get("uploaded_at", ""), "steps_done": len(steps)})
    return JSONResponse({"rfps": items})


@app.delete("/api/rfp/{rfp_id}")
async def delete_rfp(rfp_id: str):
    meta = db.get_rfp_meta(rfp_id)
    fname = meta["filename"] if meta else ""
    filepath = db.delete_rfp(rfp_id)
    if filepath and Path(filepath).exists():
        Path(filepath).unlink(missing_ok=True)
    log_activity("RFP 삭제", fname)
    return JSONResponse({"ok": True})


@app.get("/api/rfp-detail/{rfp_id}")
async def rfp_detail(rfp_id: str):
    rfp = db.get_rfp_meta(rfp_id)
    if not rfp:
        raise HTTPException(404, "RFP를 찾을 수 없습니다.")
    pipeline = db.get_pipeline(rfp_id)
    return JSONResponse({
        "rfp": {"id": rfp["id"], "filename": rfp["filename"], "text_length": rfp["text_length"], "uploaded_at": rfp.get("uploaded_at", "")},
        "steps_done": list(pipeline.get("completed_steps", {}).keys()),
        "results": pipeline.get("results", {}),
        "versions": db.list_versions(rfp_id),
        "team": db.get_team(rfp_id),
        "proposals": db.list_proposals_by_rfp(rfp_id),
    })


@app.websocket("/ws/{username}")
async def websocket_endpoint(ws: WebSocket, username: str):
    await ws_manager.connect(ws, username)
    try:
        while True:
            await ws.receive_text()  # keep alive
    except WebSocketDisconnect:
        ws_manager.disconnect(ws)
        await ws_manager.broadcast_users()


@app.on_event("startup")
async def startup():
    db.init_db()
    print(f"[Startup] DB ready. {db.count_rfps()} RFPs, {db.count_proposals()} proposals, {db.count_knowledge()} knowledge items")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
